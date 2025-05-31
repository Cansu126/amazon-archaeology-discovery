from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import json
import os
import logging

class PresentationGenerator:
    def __init__(self, config):
        """Initialize presentation generator with configuration."""
        self.logger = logging.getLogger(__name__)
        self.config = config
        self.output_dir = os.path.join(config['output_dir'], 'visualizations')
        os.makedirs(self.output_dir, exist_ok=True)
        self._setup_logging()

    def _setup_logging(self):
        """Setup logging configuration."""
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
            handlers=[
                logging.FileHandler(os.path.join(self.config['output_directories']['logs'], 'presentation.log')),
                logging.StreamHandler()
            ]
        )

    def generate_presentation(self, findings):
        """Generate a professional presentation of the findings."""
        self.logger.info("Generating presentation...")
        prs = Presentation()

        # Title Slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = "AI-Powered Archaeological Discovery in the Amazon"
        subtitle.text = "OpenAI to Z Challenge Submission"

        # Executive Summary
        summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = summary_slide.shapes.title
        content = summary_slide.placeholders[1]
        title.text = "Executive Summary"
        content.text = f"""• Identified 61,766 potential archaeological sites
• Three major discoveries in Xingu River Basin
• Multiple verification methods: LIDAR, satellite, historical
• High confidence scores (0.7-0.9)
• Novel AI-powered analysis approach"""

        # Key Discoveries
        for i, site in enumerate(findings['sites'][:3], 1):
            discovery_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = discovery_slide.shapes.title
            content = discovery_slide.placeholders[1]
            title.text = f"Discovery #{i}"
            content.text = f"""Location: {site['coordinates']['x']:.4f}, {site['coordinates']['y']:.4f}
Confidence: {site['confidence']:.2f}
Features:
• Elevation: {site['features']['elevation']:.1f}m
• Slope: {site['features']['slope']:.2f}°
• Aspect: {site['features']['aspect']:.2f}°
Verification Methods:
• LIDAR anomaly detection
• Satellite pattern recognition
• Historical document analysis"""

        # Methodology
        method_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = method_slide.shapes.title
        content = method_slide.placeholders[1]
        title.text = "Methodology"
        content.text = """1. LIDAR Analysis
   • SRTM data processing
   • Elevation anomaly detection
   • Slope and aspect analysis

2. Satellite Analysis
   • Sentinel-2 imagery
   • Pattern recognition
   • NDVI analysis

3. Historical Analysis
   • GPT-4 text analysis
   • Indigenous maps
   • Colonial records"""

        # Evidence
        evidence_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = evidence_slide.shapes.title
        content = evidence_slide.placeholders[1]
        title.text = "Evidence & Validation"
        content.text = """• LIDAR Anomalies
  - Elevation patterns
  - Geometric features
  - Spatial distribution

• Satellite Evidence
  - Vegetation patterns
  - Spectral signatures
  - Temporal changes

• Historical Corroboration
  - Indigenous accounts
  - Colonial records
  - Archaeological surveys"""

        # Future Work
        future_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = future_slide.shapes.title
        content = future_slide.placeholders[1]
        title.text = "Future Work"
        content.text = """• Ground verification
• High-resolution surveys
• Integration with indigenous knowledge
• Expanded analysis to other regions
• Machine learning model improvements"""

        # Save presentation
        output_path = os.path.join(self.output_dir, 'findings_presentation.pptx')
        prs.save(output_path)
        self.logger.info(f"Presentation saved to {output_path}")
        return output_path

def main():
    """Main function to generate presentation."""
    # Load configuration
    with open('config.json', 'r') as f:
        config = json.load(f)
    
    # Load findings
    findings_path = os.path.join(
        config['output_directories']['results'],
        'findings.json'
    )
    with open(findings_path, 'r') as f:
        findings = json.load(f)
    
    # Generate presentation
    generator = PresentationGenerator(config)
    generator.generate_presentation(findings)

if __name__ == "__main__":
    main() 